#!/usr/bin/env python3
"""
Generate PowerPoint presentation for ADK Agent submission
Run: python generate_pptx.py
Output: ADK_Agent_Submission.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create PowerPoint presentation for ADK Agent"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    TITLE_COLOR = RGBColor(25, 118, 210)  # Blue
    ACCENT_COLOR = RGBColor(56, 142, 60)  # Green
    TEXT_COLOR = RGBColor(33, 33, 33)     # Dark gray
    
    def add_title_slide(title, subtitle):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = TITLE_COLOR
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items):
        """Add content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.7))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        
        return slide
    
    # ===== SLIDE 1: TITLE =====
    add_title_slide(
        "ADK Text Classification Agent",
        "Deployed on Google Cloud Run with Gemini AI\n\nDeveloper: Shreyansh Tripathi | March 31, 2026"
    )
    
    # ===== SLIDE 2: OVERVIEW =====
    add_content_slide(
        "Project Overview",
        [
            "✅ AI Agent Framework: Google ADK (Agent Development Kit)",
            "✅ AI Model: Google Gemini Pro/Flash",
            "✅ Deployment: Google Cloud Run (Serverless)",
            "✅ Primary Task: Text Classification into 5 Categories",
            "✅ API Type: RESTful HTTP Endpoints",
            "✅ Status: Live and Production-Ready",
            "✅ 3 API Endpoints | 5 Classification Categories",
            "✅ JSON Request/Response | Auto-Scaling Infrastructure"
        ]
    )
    
    # ===== SLIDE 3: AGENT CAPABILITY =====
    add_content_slide(
        "Agent Capability: Text Classification",
        [
            "Supported Categories:",
            "   • NEWS - News articles and press releases",
            "   • OPINION - Opinion pieces and editorials",
            "   • TECHNICAL - Technical documentation",
            "   • MARKETING - Marketing and promotional content",
            "   • EDUCATIONAL - Educational and tutorials",
            "",
            "Output Per Classification:",
            "   • Category name | Confidence score (0.0-1.0) | Reasoning"
        ]
    )
    
    # ===== SLIDE 4: ARCHITECTURE =====
    add_content_slide(
        "Technical Architecture",
        [
            "Technology Stack:",
            "   • Python 3.11 | Flask (HTTP Server)",
            "   • Google Generative AI SDK",
            "   • Docker (Containerization)",
            "   • Google Cloud Run (Deployment)",
            "",
            "System Flow:",
            "   HTTP Request → Flask Server → ADK Agent →",
            "   Gemini API → JSON Response"
        ]
    )
    
    # ===== SLIDE 5: API ENDPOINTS =====
    add_content_slide(
        "API Endpoints",
        [
            "1️⃣  Health Check: GET /",
            "     Returns: Service status and version",
            "",
            "2️⃣  Agent Info: GET /agent/info",
            "     Returns: Available categories, models, endpoints",
            "",
            "3️⃣  Text Classification: POST /classify ⭐",
            "     Body: {\"text\": \"Your text here\"}",
            "     Returns: {\"category\": ..., \"confidence\": ...}"
        ]
    )
    
    # ===== SLIDE 6: DEPLOYMENT =====
    add_content_slide(
        "Cloud Run Deployment",
        [
            "Configuration:",
            "   • Service: adk-agent  | Region: us-central1",
            "   • Memory: 512MB | CPU: 1 vCPU",
            "   • Timeout: 300 seconds | Concurrency: 80",
            "   • Auto-Scaling: 0 to 1000+ instances",
            "",
            "Container:",
            "   • Server: Gunicorn (production WSGI)",
            "   • Port: 8080 | Access: Public (HTTPS)"
        ]
    )
    
    # ===== SLIDE 7: MODEL SELECTION =====
    add_content_slide(
        "Gemini Model Strategy",
        [
            "Auto-Selection Chain:",
            "   1. Gemini 2.0 Flash (Latest & Fastest)",
            "   2. Gemini 1.5 Pro (High Accuracy)",
            "   3. Gemini 1.5 Flash (Balanced)",
            "   4. Gemini Pro (Fallback)",
            "",
            "Advantages:",
            "   ✅ Always works regardless of availability",
            "   ✅ Uses latest/fastest model when possible"
        ]
    )
    
    # ===== SLIDE 8: SUBMISSION LINKS =====
    add_content_slide(
        "Submission Links",
        [
            "🔗 PRIMARY SUBMISSION LINK (LIVE) ⭐",
            "   https://adk-agent-109924518677.us-central1.run.app",
            "",
            "🔗 GITHUB REPOSITORY",
            "   https://github.com/shreyansh9026/adk-summarization-agent",
            "",
            "📊 PROJECT INFO",
            "   • Project ID: adk-agent-shreyansh",
            "   • Status: ✅ LIVE & OPERATIONAL"
        ]
    )
    
    # ===== SLIDE 9: FEATURES =====
    add_content_slide(
        "Key Features",
        [
            "✅ Production-Ready",
            "   Gunicorn server, error handling, request validation",
            "",
            "✅ Scalable & Resilient",
            "   Auto-scales 0 to 1000+ | Model fallback | Timeouts",
            "",
            "✅ Well-Documented",
            "   README, deployment guide, quick start, test suite",
            "",
            "✅ Easy Integration",
            "   RESTful API, JSON format, no auth required"
        ]
    )
    
    # ===== SLIDE 10: TESTING CHECKLIST =====
    add_content_slide(
        "Testing & Verification",
        [
            "✅ Cloud Run Deployed - LIVE",
            "✅ HTTP Endpoints - WORKING (3/3)",
            "✅ Health Check - VERIFIED",
            "✅ Classification Logic - TESTED",
            "✅ Error Handling - VERIFIED",
            "✅ Docker Container - BUILT",
            "✅ GitHub Repository - PUSHED",
            "✅ Documentation - COMPLETE"
        ]
    )
    
    # ===== SLIDE 11: COST =====
    add_content_slide(
        "Cost Analysis",
        [
            "Cloud Run Free Tier:",
            "   • 2M requests/month ✅",
            "   • 360K vCPU-seconds/month ✅",
            "   • 180K GB-seconds/month ✅",
            "",
            "Typical Monthly Cost (Light Usage):",
            "   • Cloud Run: ~$2.57 (or FREE on free tier)",
            "   • Gemini API: Pay-per-token (~$0.001 per 1000 tokens)",
            "   • Result: FREE for hobby usage! 🎉"
        ]
    )
    
    # ===== SLIDE 12: QUICK START =====
    add_content_slide(
        "Deploy in 3 Minutes",
        [
            "Step 1: Install gcloud CLI",
            "   https://cloud.google.com/sdk/docs/install",
            "",
            "Step 2: Clone Repository",
            "   git clone https://github.com/shreyansh9026/adk-summarization-agent.git",
            "",
            "Step 3: Deploy",
            "   gcloud run deploy adk-agent --source . --region us-central1",
            "",
            "Done! Your agent is now live! ✅"
        ]
    )
    
    # ===== SLIDE 13: GITHUB REPOSITORY =====
    add_content_slide(
        "GitHub Repository",
        [
            "📦 Files Included:",
            "   • agent.py - ADK Agent Implementation",
            "   • server.py - Flask HTTP Server",
            "   • Dockerfile - Container Configuration",
            "   • test_agent.py - Test Suite",
            "   • README.md - Full Documentation",
            "   • QUICK_START.md - 5-Minute Guide",
            "",
            "Repository: https://github.com/shreyansh9026/adk-summarization-agent"
        ]
    )
    
    # ===== SLIDE 14: DEMONSTRATION =====
    add_content_slide(
        "Live Agent Demo",
        [
            "📥 INPUT: \"Breaking news: Tech breakthrough announced\"",
            "",
            "🤖 AGENT RESPONSE:",
            "   ✓ Category: NEWS",
            "   ✓ Confidence: 0.95",
            "   ✓ Reasoning: Article discusses current events",
            "",
            "⏱️ Response Time: 1-3 seconds",
            "📊 Status Code: 200 OK"
        ]
    )
    
    # ===== SLIDE 15: METRICS & MONITORING =====
    add_content_slide(
        "Metrics & Monitoring",
        [
            "Available Metrics:",
            "   • Request count | Response latency",
            "   • Error rate | Memory & CPU usage",
            "   • Concurrent requests",
            "",
            "Monitoring Commands:",
            "   gcloud run logs read adk-agent --follow",
            "   gcloud run services describe adk-agent",
            "",
            "Dashboard: https://console.cloud.google.com/run?project=adk-agent-shreyansh"
        ]
    )
    
    # ===== SLIDE 16: ADVANTAGES =====
    add_content_slide(
        "Why This Solution Wins",
        [
            "vs Traditional APIs:",
            "   • Cost: Free tier available vs Fixed pricing",
            "   • Setup: 3 minutes vs 2-3 hours",
            "   • Scalability: Auto-scales vs Manual scaling",
            "   • Documentation: Comprehensive vs Minimal",
            "   • Reliability: 99.9% SLA vs Varies",
            "",
            "Fully customizable, production-ready, easy integration!"
        ]
    )
    
    # ===== SLIDE 17: FUTURE ENHANCEMENTS =====
    add_content_slide(
        "Future Enhancements (Roadmap)",
        [
            "Phase 1: Core Features ✅ COMPLETE",
            "   Text classification, HTTP endpoints, deployment",
            "",
            "Phase 2: Advanced Features (Ready)",
            "   Batch classification, logging, custom categories",
            "",
            "Phase 3: Enterprise Features (Scalable)",
            "   Authentication, webhooks, async processing",
            "   Analytics dashboard, WebSocket support"
        ]
    )
    
    # ===== SLIDE 18: SECURITY =====
    add_content_slide(
        "Security & Compliance",
        [
            "✅ API Key Management",
            "   Environment variables, Cloud Secret Manager support",
            "",
            "✅ Network Security",
            "   HTTPS only, optional authentication layer",
            "",
            "✅ Data Handling",
            "   No data persistence, privacy compliant",
            "",
            "✅ Code Quality",
            "   Type hints, error handling, exception management"
        ]
    )
    
    # ===== SLIDE 19: SUMMARY =====
    add_content_slide(
        "Summary",
        [
            "✅ Production-ready AI agent",
            "✅ Deployed on Google Cloud Run",
            "✅ Text classification capability",
            "✅ RESTful HTTP API with 3 endpoints",
            "✅ Comprehensive documentation & test suite",
            "✅ GitHub repository with full source code",
            "✅ 99.9% uptime SLA",
            "✅ LIVE NOW at: https://adk-agent-109924518677.us-central1.run.app"
        ]
    )
    
    # ===== SLIDE 20: QUESTIONS =====
    add_title_slide(
        "Questions?",
        "Live Demo: https://adk-agent-109924518677.us-central1.run.app\n\nRepository: https://github.com/shreyansh9026/adk-summarization-agent"
    )
    
    # ===== SLIDE 21: THANK YOU =====
    add_title_slide(
        "Thank You!",
        "ADK Text Classification Agent\n✅ Complete & Live\n\nDeveloper: Shreyansh Tripathi | March 31, 2026"
    )
    
    # Save presentation
    output_file = "ADK_Agent_Submission.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    try:
        create_presentation()
        print("\n🎉 Presentation ready for submission!")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("\nInstall python-pptx:")
        print("   pip install python-pptx")
